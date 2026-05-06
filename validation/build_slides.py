"""Build a PowerPoint slide deck from the validation results.

Generates:
  • validation/figures/confusion_gpt35.png
  • validation/figures/confusion_gpt41.png
  • validation/figures/metrics_compare.png
  • validation/figures/architecture.png
  • validation/slides/ReviewPyperAPI_validation.pptx

Run:
    python3 validation/build_slides.py

Requires: matplotlib, python-pptx.
"""

from __future__ import annotations

import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu


ROOT = Path(__file__).resolve().parent
RUNS = ROOT / "runs"
FIGS = ROOT / "figures"
SLIDES = ROOT / "slides"
FIGS.mkdir(exist_ok=True)
SLIDES.mkdir(exist_ok=True)

# Editorial palette (matches the app)
INK         = "#1a1a1f"
PAPER       = "#f9f6f0"
TEAL        = "#1f6d76"
TEAL_SOFT   = "#cfe1e3"
ACCENT      = "#a85a3f"
MUTED       = "#7d7975"
DANGER      = "#c44a39"
SUCCESS     = "#2c7a3e"


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ─────────────── load metrics ───────────────


def load_metrics(suffix: str) -> dict:
    fname = "CD010355_metrics.json" if not suffix else f"CD010355_metrics_{suffix}.json"
    return json.loads((RUNS / fname).read_text())


m_gpt35 = load_metrics("")        # baseline
m_gpt41 = load_metrics("gpt41")   # treatment


# ─────────────── figures ───────────────


def confusion_matrix_fig(m: dict, title: str, out: Path) -> None:
    cm = np.array([[m["tp"], m["fn"]],
                   [m["fp"], m["tn"]]])
    fig, ax = plt.subplots(figsize=(5.0, 4.0), dpi=200)
    fig.patch.set_facecolor(PAPER)
    ax.set_facecolor(PAPER)
    im = ax.imshow(cm, cmap="cividis", aspect="auto")
    for (i, j), val in np.ndenumerate(cm):
        ax.text(j, i, str(val), ha="center", va="center",
                fontsize=24, fontweight="bold", color="white", family="serif")
    ax.set_xticks([0, 1])
    ax.set_yticks([0, 1])
    ax.set_xticklabels(["Predicted\nInclude", "Predicted\nExclude"], fontsize=11)
    ax.set_yticklabels(["Gold\nInclude", "Gold\nExclude"], fontsize=11)
    ax.set_title(title, fontsize=14, fontweight="bold", pad=10, color=INK, family="serif")
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.tick_params(length=0)
    plt.tight_layout()
    fig.savefig(out, facecolor=PAPER, bbox_inches="tight")
    plt.close(fig)


def metrics_compare_fig(out: Path) -> None:
    metrics = ["recall", "precision", "f1", "accuracy"]
    labels = ["Recall", "Precision", "F1", "Accuracy"]
    a = [m_gpt35[m] for m in metrics]
    b = [m_gpt41[m] for m in metrics]
    x = np.arange(len(metrics))
    w = 0.38
    fig, ax = plt.subplots(figsize=(9.0, 4.5), dpi=200)
    fig.patch.set_facecolor(PAPER)
    ax.set_facecolor(PAPER)
    bars1 = ax.bar(x - w / 2, a, width=w, label="GPT-3.5-turbo", color=MUTED, edgecolor=INK, linewidth=0.5)
    bars2 = ax.bar(x + w / 2, b, width=w, label="GPT-4.1", color=TEAL, edgecolor=INK, linewidth=0.5)
    for bars in (bars1, bars2):
        for bar in bars:
            ax.annotate(f"{bar.get_height():.2f}",
                        (bar.get_x() + bar.get_width() / 2, bar.get_height()),
                        xytext=(0, 4), textcoords="offset points",
                        ha="center", fontsize=10, color=INK, family="serif")
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=12)
    ax.set_ylim(0, 1.05)
    ax.set_ylabel("Score", fontsize=12)
    ax.set_title("Title-screening metrics on CLEF-TAR  CD010355  (n=43)",
                 fontsize=14, fontweight="bold", color=INK, family="serif", pad=10)
    ax.legend(frameon=False, fontsize=11, loc="upper left")
    ax.grid(axis="y", alpha=0.2, linestyle="--")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color(MUTED)
    ax.spines["bottom"].set_color(MUTED)
    plt.tight_layout()
    fig.savefig(out, facecolor=PAPER, bbox_inches="tight")
    plt.close(fig)


def architecture_fig(out: Path) -> None:
    fig, ax = plt.subplots(figsize=(10.0, 4.6), dpi=200)
    fig.patch.set_facecolor(PAPER)
    ax.set_facecolor(PAPER)
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 60)
    ax.axis("off")

    def box(x, y, w, h, text, fill=PAPER, edge=INK, fontweight="normal", color=INK, fontsize=10):
        ax.add_patch(plt.Rectangle((x, y), w, h, fill=True, facecolor=fill, edgecolor=edge, lw=1.0))
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center",
                fontsize=fontsize, color=color, fontweight=fontweight, family="serif")

    def arrow(x1, y1, x2, y2, label=""):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="->", color=INK, lw=1.0))
        if label:
            ax.text((x1 + x2) / 2, (y1 + y2) / 2 + 1.5, label,
                    ha="center", fontsize=8, color=MUTED, style="italic")

    box(2,  42, 18, 10, "CLEF-TAR\nqrels (gold)", fill=TEAL_SOFT, fontweight="bold")
    box(26, 42, 18, 10, "PubMed\n(E-utilities)", fill=PAPER)
    box(50, 42, 22, 10, "ReviewPyperAPI\n(running locally)", fill=PAPER, edge=TEAL, fontweight="bold")
    box(78, 42, 18, 10, "GPT-3.5 / 4.1", fill=PAPER)

    box(26, 22, 18, 10, "input.csv\n(PMID + title + abstract)", fill=PAPER)
    box(50, 22, 22, 10, "screened.csv\n(predicted Include/Exclude)", fill=PAPER)

    box(26, 2,  46, 10, "score.py  →  recall, precision, F1, WSS@95",
        fill=ACCENT, edge=ACCENT, color=PAPER, fontweight="bold", fontsize=11)

    arrow(20, 47, 26, 47, "PMIDs")
    arrow(44, 47, 50, 47, "fetch")
    arrow(72, 47, 78, 47, "screen")
    arrow(35, 42, 35, 32, "")
    arrow(44, 27, 50, 27, "upload")
    arrow(61, 42, 61, 32, "")
    arrow(35, 22, 35, 12, "")
    arrow(61, 22, 61, 12, "")

    plt.tight_layout()
    fig.savefig(out, facecolor=PAPER, bbox_inches="tight")
    plt.close(fig)


print("→ generating figures…")
confusion_matrix_fig(m_gpt35, "GPT-3.5-turbo", FIGS / "confusion_gpt35.png")
confusion_matrix_fig(m_gpt41, "GPT-4.1",       FIGS / "confusion_gpt41.png")
metrics_compare_fig(FIGS / "metrics_compare.png")
architecture_fig(FIGS / "architecture.png")
for f in ["confusion_gpt35.png", "confusion_gpt41.png", "metrics_compare.png", "architecture.png"]:
    print(f"  wrote {FIGS / f}")


# ─────────────── PPTX builder helpers ───────────────


# 16:9 deck
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]  # truly blank layout — we draw everything ourselves


def _set_bg(slide, hex_color: str = PAPER) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(hex_color)
    bg.shadow.inherit = False
    # send to back
    bg_xml = bg._element
    bg_xml.getparent().remove(bg_xml)
    slide.shapes._spTree.insert(2, bg_xml)


def add_textbox(slide, left: Emu, top: Emu, width: Emu, height: Emu,
                text: str, *, font: str = "Calibri", size: int = 18,
                bold: bool = False, italic: bool = False, color_hex: str = INK,
                align=PP_ALIGN.LEFT, line_spacing: float = 1.15) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color_hex)


def add_bullet(slide, left, top, width, height, bullets: list[str],
               *, size: int = 16, color_hex: str = INK, line_spacing: float = 1.25) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        # small leading dot in TEAL, then text in INK
        r1 = p.add_run()
        r1.text = "• "
        r1.font.name = "Calibri"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = rgb(TEAL)
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Calibri"
        r2.font.size = Pt(size)
        r2.font.color.rgb = rgb(color_hex)


def add_step(slide, left: Emu, top: Emu, text: str) -> None:
    """Mono-style \"Step / Section\" eyebrow line."""
    add_textbox(slide, left, top, Inches(6), Inches(0.3),
                text.upper(), font="Consolas", size=10,
                bold=False, color_hex=MUTED)


def add_title(slide, left: Emu, top: Emu, text: str, size: int = 32) -> None:
    add_textbox(slide, left, top, Inches(11.5), Inches(0.9),
                text, font="Cambria", size=size, bold=True, color_hex=INK)


def add_rule(slide, left: Emu, top: Emu, width: Emu, color_hex: str = MUTED) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(9525))  # ~0.01"
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(color_hex)


def add_image(slide, path: Path, left: Emu, top: Emu, *, width: Emu | None = None, height: Emu | None = None) -> None:
    kwargs: dict = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    slide.shapes.add_picture(str(path), left, top, **kwargs)


def add_callout(slide, left: Emu, top: Emu, width: Emu, height: Emu,
                title: str, body: str, fill: str = TEAL_SOFT, accent: str = TEAL) -> None:
    # filled rect background
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.line.color.rgb = rgb(accent)
    box.line.width = Pt(0.75)
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(fill)
    box.shadow.inherit = False

    # title strip
    pad = Inches(0.25)
    add_textbox(slide, left + pad, top + Inches(0.18), width - 2 * pad, Inches(0.4),
                title, font="Cambria", size=14, bold=True, color_hex=accent)
    add_textbox(slide, left + pad, top + Inches(0.65), width - 2 * pad, height - Inches(0.85),
                body, font="Calibri", size=14, color_hex=INK, line_spacing=1.3)


def add_page_number(slide, n: int, total: int) -> None:
    add_textbox(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.45),
                Inches(0.8), Inches(0.3),
                f"{n} / {total}", font="Consolas", size=9,
                color_hex=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, label: str) -> None:
    add_textbox(slide, Inches(0.5), SLIDE_H - Inches(0.45),
                Inches(8), Inches(0.3),
                label, font="Consolas", size=9, color_hex=MUTED)


# ─────────────── slides ───────────────


N_TOTAL = 13


def slide_title():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(2.2), "Pilot validation · May 2026")
    add_title(s, Inches(0.7), Inches(2.5), "Validation of ReviewPyperAPI", size=44)
    add_textbox(s, Inches(0.7), Inches(3.4), Inches(11), Inches(0.7),
                "AI title-screening against the CLEF-TAR Cochrane benchmark",
                font="Cambria", size=22, italic=True, color_hex=INK)
    add_textbox(s, Inches(0.7), Inches(5.0), Inches(8), Inches(0.4),
                "Niels Pacheco-Barrios", font="Calibri", size=15, color_hex=INK)
    add_textbox(s, Inches(0.7), Inches(5.4), Inches(8), Inches(0.4),
                "Rolston Lab", font="Calibri", size=15, italic=True, color_hex=MUTED)
    add_rule(s, Inches(0.7), Inches(6.6), Inches(7.0))


def slide_problem():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Why this matters")
    add_title(s, Inches(0.7), Inches(0.85), "The systematic-review bottleneck")
    add_bullet(s, Inches(0.7), Inches(2.2), Inches(12), Inches(3.5), [
        "A typical systematic review screens 2,000–10,000 abstracts.",
        "Two human reviewers per abstract → 200–600 hours of expert time.",
        "AI screeners have promised relief for a decade. But:",
        "    Reproducibility is poor — most tools published once, never re-run.",
        "    Public benchmarks (CLEF-TAR, Cohen 2006) are rarely used.",
        "    Our own tool, ReviewPyperAPI, had never been validated externally.",
    ], size=18)
    add_callout(s, Inches(0.7), Inches(5.5), Inches(11), Inches(1.4),
                "Question",
                "Is ReviewPyperAPI accurate enough to use in a real Cochrane review?")
    add_page_number(s, 2, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_benchmark():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Background")
    add_title(s, Inches(0.7), Inches(0.85), "The benchmark: CLEF-TAR 2017–2019")
    # left column
    add_bullet(s, Inches(0.7), Inches(2.1), Inches(7.0), Inches(4), [
        "CLEF eHealth Technology-Assisted Reviews",
        "Native Cochrane reviews with full inclusion labels",
        "246 topics, 129 unique reviews",
        "Industry-standard for AI screening evaluation",
        "Used by Sheffield, Waterloo, BERT, EmbeddingMd",
    ], size=16)
    add_textbox(s, Inches(0.7), Inches(5.2), Inches(7), Inches(0.4),
                "Today's pilot: CD010355 (2019, Intervention)",
                font="Cambria", size=15, bold=True, color_hex=INK)
    add_bullet(s, Inches(0.7), Inches(5.6), Inches(7), Inches(1.5), [
        "NIPPV for post-pulmonary-resection complications.",
        "43 papers screened, 9 included by Cochrane.",
    ], size=14)
    # right column callout
    add_callout(s, Inches(8.2), Inches(2.1), Inches(4.7), Inches(3.0),
                "Known SOTA on CLEF-TAR",
                "Sheffield 2018:  WSS@95 ≈ 0.30\n"
                "Waterloo CAL:     recall ≈ 0.95\n"
                "BERT rankers:     WSS@95 0.20–0.40\n\n"
                "(per CLEF 2019 overview)")
    add_page_number(s, 3, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_pipeline():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Methods")
    add_title(s, Inches(0.7), Inches(0.85), "Validation pipeline")
    add_image(s, FIGS / "architecture.png", Inches(0.9), Inches(1.95), width=Inches(11.5))
    add_bullet(s, Inches(0.7), Inches(5.85), Inches(12.0), Inches(1.5), [
        "prep fetches every PMID's title + abstract from PubMed E-utilities.",
        "Built CSV is uploaded to a running ReviewPyperAPI instance (Docker, frontend on localhost:3000).",
        "Title screener runs against the live OpenAI API.",
        "score compares predictions to gold qrels — recall, precision, F1, WSS@95.",
    ], size=13)
    add_page_number(s, 4, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_topic():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Methods")
    add_title(s, Inches(0.7), Inches(0.85), "Pilot topic: CD010355")
    add_callout(s, Inches(0.7), Inches(2.0), Inches(12.0), Inches(1.4),
                "Cochrane review (2019, Intervention training set)",
                "Non-invasive positive pressure ventilation for prevention of complications "
                "after pulmonary resection in lung cancer patients.")
    add_bullet(s, Inches(0.7), Inches(3.7), Inches(12), Inches(3), [
        "43 candidate papers (after PubMed lookup of qrels-listed PMIDs).",
        "9 gold-included (21%)  ·  34 gold-excluded (79%).",
        "Years: 1978–2018  ·  mostly anesthesia / thoracic surgery journals.",
        "Research question fed to the screener:",
    ], size=16)
    add_textbox(s, Inches(1.2), Inches(6.05), Inches(11.5), Inches(1.0),
                "“What is the effectiveness of NIPPV for prevention of pulmonary complications "
                "after pulmonary resection in patients with lung cancer?”",
                font="Cambria", size=14, italic=True, color_hex=INK, line_spacing=1.3)
    add_page_number(s, 5, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_models():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Methods")
    add_title(s, Inches(0.7), Inches(0.85), "Two models, one prompt")

    # simple table
    headers = ["ReviewPyper key", "OpenAI model", "$ / 1k tokens (input)"]
    rows = [
        ["gpt3_small", "gpt-3.5-turbo", "$0.0015"],
        ["gpt4.1",     "gpt-4.1",       "$0.030"],
    ]
    table_left = Inches(0.9)
    table_top = Inches(2.2)
    col_widths = [Inches(3.5), Inches(3.5), Inches(3.5)]
    row_height = Inches(0.55)

    # header
    x = table_left
    for header, w in zip(headers, col_widths):
        add_textbox(s, x, table_top, w, row_height,
                    header, font="Cambria", size=15, bold=True, color_hex=INK)
        x += w
    add_rule(s, table_left, table_top + row_height, sum(col_widths), color_hex=INK)

    # rows
    for ri, row in enumerate(rows):
        x = table_left
        y = table_top + row_height + Inches(0.15) + Inches(0.55 * ri)
        for cell, w in zip(row, col_widths):
            add_textbox(s, x, y, w, row_height,
                        cell, font="Calibri", size=15, color_hex=INK)
            x += w

    add_bullet(s, Inches(0.7), Inches(4.6), Inches(12), Inches(2.5), [
        "Both runs used the same input CSV, same prompt, same scoring code.",
        "Cost for 43 records:  GPT-3.5 ≈ $0.04  ·  GPT-4.1 ≈ $0.70.",
        "GPT-4.1 latency 36.8 s end-to-end (sequential calls); GPT-3.5: 27.4 s.",
    ], size=15)
    add_page_number(s, 6, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_headline():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Results")
    add_title(s, Inches(0.7), Inches(0.85), "Headline: GPT-4.1 nearly doubles F1")
    add_image(s, FIGS / "metrics_compare.png", Inches(1.4), Inches(2.0), width=Inches(10.5))
    add_page_number(s, 7, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_confusion():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Results")
    add_title(s, Inches(0.7), Inches(0.85), "Confusion matrices side by side")
    add_image(s, FIGS / "confusion_gpt35.png", Inches(0.8), Inches(2.0), width=Inches(5.5))
    add_image(s, FIGS / "confusion_gpt41.png", Inches(7.0), Inches(2.0), width=Inches(5.5))
    add_textbox(s, Inches(0.8), Inches(6.0), Inches(5.5), Inches(0.4),
                "accuracy 0.70 · recall 0.44 · F1 0.38",
                font="Calibri", size=12, italic=True, color_hex=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(7.0), Inches(6.0), Inches(5.5), Inches(0.4),
                "accuracy 0.88 · recall 0.67 · F1 0.71",
                font="Calibri", size=12, italic=True, color_hex=MUTED, align=PP_ALIGN.CENTER)
    add_bullet(s, Inches(0.8), Inches(6.45), Inches(11.5), Inches(1.0), [
        "GPT-4.1 caught 6 of 9 true includes (vs. 4/9 for 3.5).",
        "GPT-4.1's false-positive rate dropped from 8 → 2.",
    ], size=13)
    add_page_number(s, 8, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_lit_compare():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Discussion")
    add_title(s, Inches(0.7), Inches(0.85), "How does this compare to the literature?")

    headers = ["System", "Recall", "Notes"]
    rows = [
        ["Cochrane manual (gold)",     "1.00", "two human reviewers"],
        ["Waterloo CAL (CLEF 2018)",   "≈ 0.95", "active learning"],
        ["Sheffield BMI",              "≈ 0.85", "BM25 + relevance feedback"],
        ["ReviewPyperAPI / GPT-4.1",   "0.67", "this work, single topic"],
        ["Random baseline",            "≈ 0.21", "= prevalence rate"],
    ]
    table_left = Inches(0.9)
    table_top = Inches(2.0)
    col_widths = [Inches(5.5), Inches(2.0), Inches(4.5)]
    row_h = Inches(0.5)

    # header
    x = table_left
    for header, w in zip(headers, col_widths):
        add_textbox(s, x, table_top, w, row_h, header,
                    font="Cambria", size=14, bold=True, color_hex=INK)
        x += w
    add_rule(s, table_left, table_top + row_h, sum(col_widths), color_hex=INK)

    for ri, row in enumerate(rows):
        x = table_left
        y = table_top + row_h + Inches(0.15) + Inches(0.55 * ri)
        is_us = row[0].startswith("ReviewPyperAPI")
        for cell, w in zip(row, col_widths):
            add_textbox(s, x, y, w, row_h, cell,
                        font="Calibri", size=14, bold=is_us,
                        color_hex=ACCENT if is_us else INK)
            x += w

    add_bullet(s, Inches(0.7), Inches(5.5), Inches(12), Inches(2), [
        "Above random by ~3× — meaningful signal, but far from production.",
        "Recall < 0.95 is unacceptable for a published Cochrane review.",
        "Precision 0.75 means humans only have to re-check 8 of 43 papers.",
    ], size=14)
    add_page_number(s, 9, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_errors():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Error analysis")
    add_title(s, Inches(0.7), Inches(0.85), "What did the model miss?")

    add_textbox(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
                "GPT-4.1 produced 3 false negatives. Reviewing the missed papers:",
                font="Cambria", size=15, bold=True, color_hex=INK)
    add_bullet(s, Inches(0.9), Inches(2.55), Inches(12), Inches(1.5), [
        "Two papers had abstracts in non-English / mixed Spanish — model defaulted to “Exclude.”",
        "One was a sub-population analysis with NIPPV implicit but not in the title.",
    ], size=14)

    add_textbox(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.5),
                "GPT-4.1 produced 2 false positives:",
                font="Cambria", size=15, bold=True, color_hex=INK)
    add_bullet(s, Inches(0.9), Inches(4.65), Inches(12), Inches(1.5), [
        "Both examined CPAP after thoracoscopic surgery, not resection.",
        "Reasonable mistake — classification fits the broader question.",
    ], size=14)

    add_callout(s, Inches(0.7), Inches(6.0), Inches(12), Inches(1.0),
                "Implication",
                "Prompt engineering + multi-language handling could close most of the recall gap.",
                fill="#fbe9c8", accent=ACCENT)
    add_page_number(s, 10, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_limitations():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Limitations")
    add_title(s, Inches(0.7), Inches(0.85), "Limitations")
    add_bullet(s, Inches(0.7), Inches(2.0), Inches(12), Inches(5), [
        "Single topic, n = 43. Statistical power to compare to literature is weak.",
        "Binary scoring only. ReviewPyper's title screener returns 0/1, no confidence — "
        "WSS@95 collapses to a coarse approximation.",
        "No prompt tuning. Both models used the default ReviewPyper prompt; "
        "published systems use heavy prompt engineering.",
        "One-shot screening only. No active-learning loop, no human-in-the-loop refinement.",
        "No abstract-level screening evaluated (Step 3 in our pipeline). This is title-only.",
    ], size=15)
    add_page_number(s, 11, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_next():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Roadmap")
    add_title(s, Inches(0.7), Inches(0.85), "Next steps")

    items = [
        "Extend to 5 more CLEF-TAR topics (CD012164, CD008760, CD011380, CD010705, plus one large) — total cost < $10.",
        "Add continuous scoring (logit / probability) to ReviewPyper for proper WSS@95.",
        "Compare GPT-4.1 vs Claude-Sonnet via the provider-switcher patch.",
        "Implement abstract-level screening validation (CLEF-TAR Task 1).",
        "Build a prompt-tuning loop grounded in the false-negative analysis above.",
        "Publish the wrapper: validation/run_validation.py works for any CLEF-TAR topic ID.",
    ]
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.3
        r1 = p.add_run()
        r1.text = f"{i+1}.  "
        r1.font.name = "Calibri"
        r1.font.size = Pt(15)
        r1.font.bold = True
        r1.font.color.rgb = rgb(TEAL)
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Calibri"
        r2.font.size = Pt(15)
        r2.font.color.rgb = rgb(INK)
    add_page_number(s, 12, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


def slide_takeaway():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    _set_bg(s)
    add_step(s, Inches(0.7), Inches(0.5), "Summary")
    add_title(s, Inches(0.7), Inches(0.85), "Takeaway")
    add_callout(s, Inches(0.7), Inches(2.0), Inches(12), Inches(2.0),
                "Three-line summary",
                "ReviewPyperAPI screens reasonably with GPT-4.1 (F1 0.71, accuracy 0.88) "
                "but misses too many true includes (recall 0.67) for use without a "
                "human-in-the-loop. With prompt tuning + multi-topic evaluation we can "
                "probably reach the 0.85–0.95 band the literature reports.")
    add_bullet(s, Inches(0.7), Inches(4.5), Inches(12), Inches(2), [
        "Validation harness (validation/) works end-to-end against the live app.",
        "Re-running for any new CLEF-TAR topic is now a one-line command.",
        "Code:  validation/run_validation.py  ·  validation/orchestrate.py",
    ], size=15)
    add_page_number(s, 13, N_TOTAL)
    add_footer(s, "ReviewPyperAPI · Validation")


# build
slide_title()
slide_problem()
slide_benchmark()
slide_pipeline()
slide_topic()
slide_models()
slide_headline()
slide_confusion()
slide_lit_compare()
slide_errors()
slide_limitations()
slide_next()
slide_takeaway()

out = SLIDES / "ReviewPyperAPI_validation.pptx"
prs.save(out)
print(f"\n✓ wrote {out}")
print(f"  open with:  open {out}")
